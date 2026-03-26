import sys
import os

try:
    from PyPDF2 import PdfReader
    def extract_pdf(path):
        text = ""
        try:
            reader = PdfReader(path)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        except Exception as e:
            text = f"Error reading PDF: {e}"
        return text
except ImportError:
    def extract_pdf(path): return "PyPDF2 not installed"

try:
    import docx
    def extract_docx(path):
        text = ""
        try:
            doc = docx.Document(path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        except Exception as e:
            text = f"Error reading DOCX: {e}"
        return text
except ImportError:
    def extract_docx(path): return "python-docx not installed"

try:
    from pptx import Presentation
    def extract_pptx(path):
        text = ""
        try:
            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            text = f"Error reading PPTX: {e}"
        return text
except ImportError:
    def extract_pptx(path): return "python-pptx not installed"

def main():
    if len(sys.argv) < 2:
        return
    path = sys.argv[1]
    ext = os.path.splitext(path)[1].lower()
    
    if ext == '.pdf':
        print(extract_pdf(path))
    elif ext == '.docx':
        print(extract_docx(path))
    elif ext == '.pptx':
        print(extract_pptx(path))
    else:
        print("Unsupported format")

if __name__ == "__main__":
    main()
